"""
Offline Document Parser Service

Parses various document types locally without external API calls.
Designed for air-gapped environments.

Supported formats:
- PDF (via PyMuPDF/fitz with fallback to PyPDF2)
- Word (.docx)
- PowerPoint (.pptx)
- Excel (.xlsx)
- Text files (.txt, .md, .json, .yaml, .xml)
- Images (OCR via tesseract if available)

All processing is done locally - no external API calls.
"""

import base64
import io
import json
import re
import os
from typing import Dict, List, Any, Optional, Tuple
from dataclasses import dataclass, field
from enum import Enum
import tempfile

from backend.core.logging import get_logger
from backend.services.tesseract_utils import configure_pytesseract

logger = get_logger(__name__)


class DocumentType(str, Enum):
    """Supported document types."""
    PDF = "pdf"
    WORD = "docx"
    POWERPOINT = "pptx"
    EXCEL = "xlsx"
    TEXT = "text"
    MARKDOWN = "markdown"
    JSON = "json"
    YAML = "yaml"
    XML = "xml"
    HTML = "html"
    IMAGE = "image"
    CODE = "code"
    OPENAPI = "openapi"
    UNKNOWN = "unknown"


@dataclass
class DocumentSection:
    """A section of extracted document content."""
    title: str
    content: str
    page_number: Optional[int] = None
    section_type: str = "text"  # text, heading, table, code, list
    importance_score: float = 0.5  # 0-1, higher = more security relevant
    keywords_found: List[str] = field(default_factory=list)


@dataclass
class ParsedDocument:
    """Result of parsing a document."""
    filename: str
    document_type: DocumentType
    total_chars: int
    total_pages: Optional[int] = None

    # Full extracted text
    full_text: str = ""

    # Structured sections (for smarter truncation)
    sections: List[DocumentSection] = field(default_factory=list)

    # Metadata
    metadata: Dict[str, Any] = field(default_factory=dict)

    # Security-relevant excerpts (prioritized content)
    security_excerpts: List[str] = field(default_factory=list)

    # API endpoints found (for OpenAPI/Swagger docs)
    api_endpoints: List[Dict[str, Any]] = field(default_factory=list)

    # Errors during parsing
    parse_errors: List[str] = field(default_factory=list)

    def get_prioritized_content(self, max_chars: int = 50000) -> str:
        """
        Get document content prioritized by security relevance.
        Returns up to max_chars of the most important content.
        """
        result_parts = []
        current_chars = 0

        # 1. Always include security excerpts first (highest priority)
        if self.security_excerpts:
            excerpt_text = "\n\n## Security-Relevant Sections\n\n" + "\n\n".join(self.security_excerpts)
            if current_chars + len(excerpt_text) <= max_chars:
                result_parts.append(excerpt_text)
                current_chars += len(excerpt_text)

        # 2. Include API endpoints if found
        if self.api_endpoints:
            api_text = "\n\n## API Endpoints\n\n"
            for ep in self.api_endpoints[:50]:  # Top 50 endpoints
                api_text += f"- {ep.get('method', 'GET')} {ep.get('path', '/')} - {ep.get('summary', '')}\n"
            if current_chars + len(api_text) <= max_chars:
                result_parts.append(api_text)
                current_chars += len(api_text)

        # 3. Include high-importance sections
        high_importance = sorted(
            [s for s in self.sections if s.importance_score >= 0.7],
            key=lambda s: s.importance_score,
            reverse=True
        )
        for section in high_importance:
            section_text = f"\n### {section.title}\n{section.content}"
            if current_chars + len(section_text) <= max_chars:
                result_parts.append(section_text)
                current_chars += len(section_text)

        # 4. Fill remaining space with full text
        remaining = max_chars - current_chars
        if remaining > 1000 and self.full_text:
            # Take from beginning and end for context
            half = remaining // 2
            beginning = self.full_text[:half]
            ending = self.full_text[-half:] if len(self.full_text) > half else ""
            if ending and beginning != ending:
                result_parts.append(f"\n\n## Document Content\n\n{beginning}\n\n[...]\n\n{ending}")
            else:
                result_parts.append(f"\n\n## Document Content\n\n{beginning}")

        return "".join(result_parts)

    def to_dict(self) -> Dict[str, Any]:
        return {
            "filename": self.filename,
            "document_type": self.document_type.value,
            "total_chars": self.total_chars,
            "total_pages": self.total_pages,
            "sections_count": len(self.sections),
            "security_excerpts_count": len(self.security_excerpts),
            "api_endpoints_count": len(self.api_endpoints),
            "metadata": self.metadata,
            "parse_errors": self.parse_errors,
        }

    def get_smart_summary(self, max_chars: int = 10000) -> str:
        """
        Generate a smart summary for huge documents.
        Extracts key information without losing critical security context.
        """
        parts = []
        current = 0

        # Header with document stats
        header = f"## {self.filename} ({self.document_type.value})\n"
        header += f"Total: {self.total_chars:,} chars, {self.total_pages or 'N/A'} pages, "
        header += f"{len(self.sections)} sections, {len(self.security_excerpts)} security excerpts\n\n"
        parts.append(header)
        current += len(header)

        # All security excerpts (condensed)
        if self.security_excerpts:
            parts.append("### Key Security Content\n")
            current += 25
            for excerpt in self.security_excerpts[:15]:
                condensed = excerpt[:300].replace("\n", " ").strip()
                if current + len(condensed) + 5 > max_chars:
                    break
                parts.append(f"- {condensed}\n")
                current += len(condensed) + 5

        # API endpoints (critical for security analysis)
        if self.api_endpoints and current < max_chars - 500:
            parts.append("\n### API Endpoints\n")
            current += 20
            for ep in self.api_endpoints[:30]:
                line = f"- {ep.get('method', 'GET')} {ep.get('path', '/')}"
                if ep.get('security'):
                    line += " [AUTH]"
                line += f": {ep.get('summary', '')[:50]}\n"
                if current + len(line) > max_chars:
                    break
                parts.append(line)
                current += len(line)

        # High-importance section titles (for context)
        high_sections = [s for s in self.sections if s.importance_score >= 0.6]
        if high_sections and current < max_chars - 200:
            parts.append("\n### Important Sections\n")
            for section in high_sections[:10]:
                line = f"- {section.title} (score: {section.importance_score:.1f})\n"
                if current + len(line) > max_chars:
                    break
                parts.append(line)
                current += len(line)

        return "".join(parts)


# Security-relevant keywords for importance scoring
SECURITY_KEYWORDS = [
    # Authentication
    "authentication", "auth", "login", "password", "credential", "token", "jwt",
    "oauth", "saml", "sso", "session", "cookie", "bearer", "api key", "secret",

    # Authorization
    "authorization", "permission", "role", "access control", "rbac", "acl",
    "privilege", "admin", "root", "sudo", "elevated",

    # Vulnerabilities
    "vulnerability", "exploit", "injection", "xss", "csrf", "sqli", "rce",
    "ssrf", "xxe", "idor", "insecure", "unsafe", "dangerous", "critical",
    "cve", "cwe", "owasp",

    # Cryptography
    "encrypt", "decrypt", "hash", "salt", "cipher", "ssl", "tls", "certificate",
    "private key", "public key", "rsa", "aes", "hmac",

    # Input/Output
    "input", "output", "sanitize", "validate", "escape", "encode", "decode",
    "filter", "whitelist", "blacklist", "regex",

    # Network
    "endpoint", "api", "rest", "graphql", "websocket", "http", "https",
    "port", "firewall", "proxy", "cors", "header",

    # Data
    "database", "sql", "query", "sensitive", "pii", "personal", "confidential",
    "credit card", "ssn", "health", "financial",

    # File operations
    "file", "upload", "download", "path", "directory", "traversal",

    # Error handling
    "error", "exception", "stack trace", "debug", "logging",
]


class DocumentParserService:
    """
    Offline document parsing service.
    All processing done locally - no external API calls.
    """

    def __init__(self):
        self._check_dependencies()

    def _check_dependencies(self):
        """Check which optional dependencies are available."""
        self._has_pymupdf = False
        self._has_pypdf2 = False
        self._has_docx = False
        self._has_pptx = False
        self._has_openpyxl = False
        self._has_tesseract = False
        self._has_pytesseract = False
        self._has_pil = False
        self._has_yaml = False

        try:
            import fitz  # PyMuPDF
            self._has_pymupdf = True
            logger.info("PyMuPDF available for PDF parsing")
        except ImportError:
            logger.warning("PyMuPDF not available, will try PyPDF2")

        try:
            import PyPDF2
            self._has_pypdf2 = True
            logger.info("PyPDF2 available for PDF parsing")
        except ImportError:
            logger.warning("PyPDF2 not available")

        try:
            import docx
            self._has_docx = True
        except ImportError:
            logger.warning("python-docx not available")

        try:
            import pptx
            self._has_pptx = True
        except ImportError:
            logger.warning("python-pptx not available")

        try:
            import openpyxl
            self._has_openpyxl = True
        except ImportError:
            logger.warning("openpyxl not available")

        try:
            from PIL import Image
            self._has_pil = True
        except ImportError:
            logger.warning("Pillow not available")

        tesseract_cmd = None
        try:
            import pytesseract
            self._has_pytesseract = True
            tesseract_cmd = configure_pytesseract()
        except ImportError:
            logger.info("pytesseract not available")

        try:
            import yaml
            self._has_yaml = True
        except ImportError:
            logger.warning("PyYAML not available")

        # Check for tesseract (optional OCR)
        if tesseract_cmd and self._has_pytesseract and self._has_pil:
            self._has_tesseract = True
            logger.info("Tesseract OCR available")
        else:
            logger.info("Tesseract OCR not available (optional)")

    def detect_document_type(
        self,
        filename: str,
        content_type: Optional[str] = None,
        content: Optional[bytes] = None,
    ) -> DocumentType:
        """Detect document type from filename, content-type, or content."""
        filename_lower = filename.lower()

        # Check by extension first
        if filename_lower.endswith(".pdf"):
            return DocumentType.PDF
        elif filename_lower.endswith(".docx"):
            return DocumentType.WORD
        elif filename_lower.endswith(".doc"):
            return DocumentType.WORD  # Old Word format
        elif filename_lower.endswith(".pptx") or filename_lower.endswith(".ppt"):
            return DocumentType.POWERPOINT
        elif filename_lower.endswith(".xlsx") or filename_lower.endswith(".xls"):
            return DocumentType.EXCEL
        elif filename_lower.endswith(".md"):
            return DocumentType.MARKDOWN
        elif filename_lower.endswith(".json"):
            return DocumentType.JSON
        elif filename_lower.endswith((".yaml", ".yml")):
            return DocumentType.YAML
        elif filename_lower.endswith(".xml"):
            return DocumentType.XML
        elif filename_lower.endswith((".html", ".htm")):
            return DocumentType.HTML
        elif filename_lower.endswith(".txt"):
            return DocumentType.TEXT
        elif filename_lower.endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff")):
            return DocumentType.IMAGE
        elif filename_lower.endswith((".py", ".js", ".ts", ".java", ".c", ".cpp", ".go", ".rs", ".rb", ".php")):
            return DocumentType.CODE

        # Check content-type
        if content_type:
            if "pdf" in content_type:
                return DocumentType.PDF
            elif "word" in content_type or "docx" in content_type:
                return DocumentType.WORD
            elif "powerpoint" in content_type or "pptx" in content_type:
                return DocumentType.POWERPOINT
            elif "excel" in content_type or "xlsx" in content_type or "spreadsheet" in content_type:
                return DocumentType.EXCEL
            elif "json" in content_type:
                return DocumentType.JSON
            elif "yaml" in content_type:
                return DocumentType.YAML
            elif "xml" in content_type:
                return DocumentType.XML
            elif "html" in content_type:
                return DocumentType.HTML
            elif "text" in content_type:
                return DocumentType.TEXT
            elif "image" in content_type:
                return DocumentType.IMAGE

        # Try to detect OpenAPI spec
        if content:
            try:
                text = content.decode("utf-8", errors="ignore")[:2000]
                if "openapi" in text.lower() or "swagger" in text.lower():
                    return DocumentType.OPENAPI
            except Exception:
                pass

        return DocumentType.UNKNOWN

    def parse_document(
        self,
        filename: str,
        content_base64: str,
        content_type: Optional[str] = None,
        description: Optional[str] = None,
        ocr_languages: Optional[str] = None,
    ) -> ParsedDocument:
        """
        Parse a document and extract text content.

        Args:
            filename: Original filename
            content_base64: Base64 encoded file content
            content_type: MIME type if known
            description: User's description of the document

        Returns:
            ParsedDocument with extracted content
        """
        # Decode base64
        try:
            content = base64.b64decode(content_base64)
        except Exception as e:
            return ParsedDocument(
                filename=filename,
                document_type=DocumentType.UNKNOWN,
                total_chars=0,
                parse_errors=[f"Failed to decode base64: {e}"],
            )

        # Detect document type
        doc_type = self.detect_document_type(filename, content_type, content)

        # Parse based on type
        if doc_type == DocumentType.PDF:
            return self._parse_pdf(filename, content, description, ocr_languages)
        elif doc_type == DocumentType.WORD:
            return self._parse_word(filename, content, description)
        elif doc_type == DocumentType.POWERPOINT:
            return self._parse_powerpoint(filename, content, description)
        elif doc_type == DocumentType.EXCEL:
            return self._parse_excel(filename, content, description)
        elif doc_type in [DocumentType.JSON, DocumentType.OPENAPI]:
            return self._parse_json(filename, content, description)
        elif doc_type == DocumentType.YAML:
            return self._parse_yaml(filename, content, description)
        elif doc_type == DocumentType.XML:
            return self._parse_xml(filename, content, description)
        elif doc_type == DocumentType.HTML:
            return self._parse_html(filename, content, description)
        elif doc_type in [DocumentType.TEXT, DocumentType.MARKDOWN, DocumentType.CODE]:
            return self._parse_text(filename, content, description, doc_type)
        elif doc_type == DocumentType.IMAGE:
            return self._parse_image(filename, content, description, ocr_languages)
        else:
            # Try as text
            return self._parse_text(filename, content, description, doc_type)

    def _parse_pdf(
        self,
        filename: str,
        content: bytes,
        description: Optional[str] = None,
        ocr_languages: Optional[str] = None,
    ) -> ParsedDocument:
        """Parse PDF using PyMuPDF (preferred) or PyPDF2 (fallback)."""
        result = ParsedDocument(
            filename=filename,
            document_type=DocumentType.PDF,
            total_chars=0,
        )

        # Try PyMuPDF first (better quality)
        if self._has_pymupdf:
            try:
                import fitz

                pdf = fitz.open(stream=content, filetype="pdf")
                result.total_pages = len(pdf)
                result.metadata = dict(pdf.metadata) if pdf.metadata else {}

                all_text = []
                ocr_used = False
                for page_num, page in enumerate(pdf):
                    text = page.get_text() or ""
                    if not text.strip() and self._has_tesseract:
                        ocr_text = self._ocr_pdf_page(page, ocr_languages)
                        if ocr_text.strip():
                            text = ocr_text
                            ocr_used = True
                    all_text.append(text)

                    # Create section for each page
                    if text.strip():
                        section = self._create_section(
                            title=f"Page {page_num + 1}",
                            content=text,
                            page_number=page_num + 1,
                        )
                        result.sections.append(section)

                pdf.close()
                result.full_text = "\n\n".join(all_text)
                result.total_chars = len(result.full_text)
                if ocr_used:
                    result.metadata["ocr_used"] = True

            except Exception as e:
                result.parse_errors.append(f"PyMuPDF error: {e}")
                # Fall through to PyPDF2

        # Fallback to PyPDF2
        if not result.full_text and self._has_pypdf2:
            try:
                import PyPDF2

                pdf_file = io.BytesIO(content)
                reader = PyPDF2.PdfReader(pdf_file)
                result.total_pages = len(reader.pages)

                if reader.metadata:
                    result.metadata = {
                        k: str(v) for k, v in reader.metadata.items()
                        if v is not None
                    }

                all_text = []
                for page_num, page in enumerate(reader.pages):
                    text = page.extract_text() or ""
                    all_text.append(text)

                    if text.strip():
                        section = self._create_section(
                            title=f"Page {page_num + 1}",
                            content=text,
                            page_number=page_num + 1,
                        )
                        result.sections.append(section)

                result.full_text = "\n\n".join(all_text)
                result.total_chars = len(result.full_text)

            except Exception as e:
                result.parse_errors.append(f"PyPDF2 error: {e}")

        # Extract security excerpts
        result.security_excerpts = self._extract_security_excerpts(result.full_text)

        # Add description context if provided
        if description:
            result.metadata["user_description"] = description

        return result

    def _parse_image(
        self,
        filename: str,
        content: bytes,
        description: Optional[str] = None,
        ocr_languages: Optional[str] = None,
    ) -> ParsedDocument:
        """Parse image using OCR if available."""
        result = ParsedDocument(
            filename=filename,
            document_type=DocumentType.IMAGE,
            total_chars=0,
        )

        if not self._has_tesseract:
            result.parse_errors.append("OCR unavailable (tesseract or Pillow missing)")
            return result

        try:
            from PIL import Image
            import pytesseract

            image = Image.open(io.BytesIO(content))
            lang = ocr_languages or "eng"
            config = "--oem 3 --psm 6"
            text = pytesseract.image_to_string(image, lang=lang, config=config)
            result.full_text = text
            result.total_chars = len(text)
            result.metadata["ocr_used"] = True
            result.security_excerpts = self._extract_security_excerpts(text)
        except Exception as e:
            result.parse_errors.append(f"OCR error: {e}")

        if description:
            result.metadata["user_description"] = description

        return result

    def _ocr_pdf_page(self, page: "fitz.Page", ocr_languages: Optional[str]) -> str:
        """OCR a PDF page using pytesseract."""
        if not self._has_tesseract:
            return ""
        try:
            from PIL import Image
            import pytesseract

            pix = page.get_pixmap(dpi=300)
            mode = "RGB" if pix.alpha == 0 else "RGBA"
            image = Image.frombytes(mode, [pix.width, pix.height], pix.samples)
            if mode == "RGBA":
                image = image.convert("RGB")
            lang = ocr_languages or "eng"
            config = "--oem 3 --psm 6"
            return pytesseract.image_to_string(image, lang=lang, config=config) or ""
        except Exception as e:
            logger.warning(f"OCR failed for PDF page: {e}")
            return ""

    def _parse_word(
        self,
        filename: str,
        content: bytes,
        description: Optional[str] = None,
    ) -> ParsedDocument:
        """Parse Word document (.docx)."""
        result = ParsedDocument(
            filename=filename,
            document_type=DocumentType.WORD,
            total_chars=0,
        )

        if not self._has_docx:
            result.parse_errors.append("python-docx not available")
            # Try as raw text
            result.full_text = content.decode("utf-8", errors="ignore")
            result.total_chars = len(result.full_text)
            return result

        try:
            from docx import Document

            doc_file = io.BytesIO(content)
            doc = Document(doc_file)

            # Extract core properties
            if doc.core_properties:
                result.metadata = {
                    "title": doc.core_properties.title or "",
                    "author": doc.core_properties.author or "",
                    "subject": doc.core_properties.subject or "",
                }

            all_text = []
            current_section_title = "Document Content"
            current_section_content = []

            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                # Check if this is a heading
                if para.style and "Heading" in para.style.name:
                    # Save previous section
                    if current_section_content:
                        section = self._create_section(
                            title=current_section_title,
                            content="\n".join(current_section_content),
                        )
                        result.sections.append(section)

                    current_section_title = text
                    current_section_content = []
                else:
                    current_section_content.append(text)
                    all_text.append(text)

            # Save last section
            if current_section_content:
                section = self._create_section(
                    title=current_section_title,
                    content="\n".join(current_section_content),
                )
                result.sections.append(section)

            # Extract tables
            for table_idx, table in enumerate(doc.tables):
                table_text = []
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    table_text.append(" | ".join(row_text))

                if table_text:
                    table_content = "\n".join(table_text)
                    all_text.append(table_content)
                    section = self._create_section(
                        title=f"Table {table_idx + 1}",
                        content=table_content,
                        section_type="table",
                    )
                    result.sections.append(section)

            result.full_text = "\n\n".join(all_text)
            result.total_chars = len(result.full_text)
            result.security_excerpts = self._extract_security_excerpts(result.full_text)

        except Exception as e:
            result.parse_errors.append(f"Word parsing error: {e}")
            # Fallback to raw text
            result.full_text = content.decode("utf-8", errors="ignore")
            result.total_chars = len(result.full_text)

        if description:
            result.metadata["user_description"] = description

        return result

    def _parse_powerpoint(
        self,
        filename: str,
        content: bytes,
        description: Optional[str] = None,
    ) -> ParsedDocument:
        """Parse PowerPoint (.pptx)."""
        result = ParsedDocument(
            filename=filename,
            document_type=DocumentType.POWERPOINT,
            total_chars=0,
        )

        if not self._has_pptx:
            result.parse_errors.append("python-pptx not available")
            return result

        try:
            from pptx import Presentation

            ppt_file = io.BytesIO(content)
            prs = Presentation(ppt_file)

            result.total_pages = len(prs.slides)
            all_text = []

            for slide_num, slide in enumerate(prs.slides):
                slide_text = []

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)

                    # Extract table content
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = [cell.text.strip() for cell in row.cells]
                            slide_text.append(" | ".join(row_text))

                if slide_text:
                    content_text = "\n".join(slide_text)
                    all_text.append(content_text)
                    section = self._create_section(
                        title=f"Slide {slide_num + 1}",
                        content=content_text,
                        page_number=slide_num + 1,
                    )
                    result.sections.append(section)

            result.full_text = "\n\n".join(all_text)
            result.total_chars = len(result.full_text)
            result.security_excerpts = self._extract_security_excerpts(result.full_text)

        except Exception as e:
            result.parse_errors.append(f"PowerPoint parsing error: {e}")

        if description:
            result.metadata["user_description"] = description

        return result

    def _parse_excel(
        self,
        filename: str,
        content: bytes,
        description: Optional[str] = None,
    ) -> ParsedDocument:
        """Parse Excel (.xlsx)."""
        result = ParsedDocument(
            filename=filename,
            document_type=DocumentType.EXCEL,
            total_chars=0,
        )

        if not self._has_openpyxl:
            result.parse_errors.append("openpyxl not available")
            return result

        try:
            import openpyxl

            xlsx_file = io.BytesIO(content)
            wb = openpyxl.load_workbook(xlsx_file, read_only=True, data_only=True)

            all_text = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_text = []

                for row in sheet.iter_rows(max_row=1000):  # Limit rows
                    row_values = []
                    for cell in row:
                        if cell.value is not None:
                            row_values.append(str(cell.value))
                    if row_values:
                        sheet_text.append(" | ".join(row_values))

                if sheet_text:
                    content_text = "\n".join(sheet_text)
                    all_text.append(content_text)
                    section = self._create_section(
                        title=f"Sheet: {sheet_name}",
                        content=content_text,
                        section_type="table",
                    )
                    result.sections.append(section)

            wb.close()
            result.full_text = "\n\n".join(all_text)
            result.total_chars = len(result.full_text)
            result.security_excerpts = self._extract_security_excerpts(result.full_text)

        except Exception as e:
            result.parse_errors.append(f"Excel parsing error: {e}")

        if description:
            result.metadata["user_description"] = description

        return result

    def _parse_json(
        self,
        filename: str,
        content: bytes,
        description: Optional[str] = None,
    ) -> ParsedDocument:
        """Parse JSON, including OpenAPI/Swagger specs."""
        result = ParsedDocument(
            filename=filename,
            document_type=DocumentType.JSON,
            total_chars=0,
        )

        try:
            text = content.decode("utf-8", errors="ignore")
            data = json.loads(text)

            # Check if OpenAPI spec
            if isinstance(data, dict) and ("openapi" in data or "swagger" in data):
                result.document_type = DocumentType.OPENAPI
                result = self._parse_openapi(result, data)
            else:
                result.full_text = json.dumps(data, indent=2)

            result.total_chars = len(result.full_text or text)
            result.security_excerpts = self._extract_security_excerpts(result.full_text or text)

        except json.JSONDecodeError as e:
            result.parse_errors.append(f"JSON parse error: {e}")
            result.full_text = content.decode("utf-8", errors="ignore")
            result.total_chars = len(result.full_text)

        if description:
            result.metadata["user_description"] = description

        return result

    def _parse_openapi(
        self,
        result: ParsedDocument,
        data: Dict[str, Any],
    ) -> ParsedDocument:
        """Extract API information from OpenAPI/Swagger spec."""
        result.metadata["openapi_version"] = data.get("openapi", data.get("swagger", "unknown"))

        info = data.get("info", {})
        result.metadata["api_title"] = info.get("title", "")
        result.metadata["api_description"] = info.get("description", "")
        result.metadata["api_version"] = info.get("version", "")

        # Extract servers
        servers = data.get("servers", [])
        result.metadata["servers"] = [s.get("url", "") for s in servers]

        # Extract paths/endpoints
        paths = data.get("paths", {})
        all_text_parts = []

        for path, methods in paths.items():
            if not isinstance(methods, dict):
                continue

            for method, details in methods.items():
                if method.startswith("x-") or not isinstance(details, dict):
                    continue

                endpoint = {
                    "path": path,
                    "method": method.upper(),
                    "summary": details.get("summary", ""),
                    "description": details.get("description", ""),
                    "tags": details.get("tags", []),
                    "security": details.get("security", []),
                    "parameters": [],
                }

                # Extract parameters
                for param in details.get("parameters", []):
                    if isinstance(param, dict):
                        endpoint["parameters"].append({
                            "name": param.get("name", ""),
                            "in": param.get("in", ""),
                            "required": param.get("required", False),
                            "type": param.get("schema", {}).get("type", ""),
                        })

                result.api_endpoints.append(endpoint)

                # Add to text for searching
                all_text_parts.append(
                    f"{method.upper()} {path}: {details.get('summary', '')} - {details.get('description', '')}"
                )

        # Extract security schemes
        security_schemes = data.get("components", {}).get("securitySchemes", {})
        if security_schemes:
            result.metadata["security_schemes"] = list(security_schemes.keys())

            security_text = "\n## Security Schemes\n"
            for name, scheme in security_schemes.items():
                security_text += f"- {name}: {scheme.get('type', '')} ({scheme.get('scheme', '')})\n"
            all_text_parts.append(security_text)

        result.full_text = "\n".join(all_text_parts)

        # Create section for API overview
        section = self._create_section(
            title="API Overview",
            content=f"API: {info.get('title', 'Unknown')}\n"
                   f"Version: {info.get('version', '')}\n"
                   f"Endpoints: {len(result.api_endpoints)}\n"
                   f"Description: {info.get('description', '')}",
            importance_score=0.9,
        )
        result.sections.append(section)

        return result

    def _parse_yaml(
        self,
        filename: str,
        content: bytes,
        description: Optional[str] = None,
    ) -> ParsedDocument:
        """Parse YAML files."""
        result = ParsedDocument(
            filename=filename,
            document_type=DocumentType.YAML,
            total_chars=0,
        )

        text = content.decode("utf-8", errors="ignore")
        result.full_text = text
        result.total_chars = len(text)

        if self._has_yaml:
            try:
                import yaml
                data = yaml.safe_load(text)

                # Check if OpenAPI spec
                if isinstance(data, dict) and ("openapi" in data or "swagger" in data):
                    result.document_type = DocumentType.OPENAPI
                    result = self._parse_openapi(result, data)

            except Exception as e:
                result.parse_errors.append(f"YAML parse error: {e}")

        result.security_excerpts = self._extract_security_excerpts(text)

        if description:
            result.metadata["user_description"] = description

        return result

    def _parse_xml(
        self,
        filename: str,
        content: bytes,
        description: Optional[str] = None,
    ) -> ParsedDocument:
        """Parse XML files safely (no XXE)."""
        result = ParsedDocument(
            filename=filename,
            document_type=DocumentType.XML,
            total_chars=0,
        )

        try:
            import defusedxml.ElementTree as ET

            text = content.decode("utf-8", errors="ignore")
            result.full_text = text
            result.total_chars = len(text)

            # Parse XML safely
            root = ET.fromstring(content)

            # Extract text content
            def extract_text(elem, depth=0):
                texts = []
                if elem.text and elem.text.strip():
                    texts.append(elem.text.strip())
                for child in elem:
                    texts.extend(extract_text(child, depth + 1))
                if elem.tail and elem.tail.strip():
                    texts.append(elem.tail.strip())
                return texts

            all_text = extract_text(root)
            result.full_text = "\n".join(all_text) if all_text else text
            result.security_excerpts = self._extract_security_excerpts(result.full_text)

        except Exception as e:
            result.parse_errors.append(f"XML parse error: {e}")
            result.full_text = content.decode("utf-8", errors="ignore")
            result.total_chars = len(result.full_text)

        if description:
            result.metadata["user_description"] = description

        return result

    def _parse_html(
        self,
        filename: str,
        content: bytes,
        description: Optional[str] = None,
    ) -> ParsedDocument:
        """Parse HTML files."""
        result = ParsedDocument(
            filename=filename,
            document_type=DocumentType.HTML,
            total_chars=0,
        )

        try:
            from bs4 import BeautifulSoup

            text = content.decode("utf-8", errors="ignore")
            soup = BeautifulSoup(text, "html.parser")

            # Extract text
            result.full_text = soup.get_text(separator="\n", strip=True)
            result.total_chars = len(result.full_text)

            # Extract title
            if soup.title:
                result.metadata["title"] = soup.title.string

            # Extract headings as sections
            for heading in soup.find_all(["h1", "h2", "h3"]):
                section_title = heading.get_text(strip=True)
                section_content = []

                # Get content until next heading
                for sibling in heading.next_siblings:
                    if sibling.name in ["h1", "h2", "h3"]:
                        break
                    if hasattr(sibling, "get_text"):
                        section_content.append(sibling.get_text(strip=True))

                if section_content:
                    section = self._create_section(
                        title=section_title,
                        content="\n".join(section_content),
                    )
                    result.sections.append(section)

            result.security_excerpts = self._extract_security_excerpts(result.full_text)

        except Exception as e:
            result.parse_errors.append(f"HTML parse error: {e}")
            result.full_text = content.decode("utf-8", errors="ignore")
            result.total_chars = len(result.full_text)

        if description:
            result.metadata["user_description"] = description

        return result

    def _parse_text(
        self,
        filename: str,
        content: bytes,
        description: Optional[str] = None,
        doc_type: DocumentType = DocumentType.TEXT,
    ) -> ParsedDocument:
        """Parse plain text files."""
        result = ParsedDocument(
            filename=filename,
            document_type=doc_type,
            total_chars=0,
        )

        # Try different encodings
        text = None
        for encoding in ["utf-8", "utf-16", "latin-1", "cp1252"]:
            try:
                text = content.decode(encoding)
                result.metadata["encoding"] = encoding
                break
            except Exception:
                continue

        if text is None:
            text = content.decode("utf-8", errors="ignore")
            result.metadata["encoding"] = "utf-8 (with errors ignored)"

        result.full_text = text
        result.total_chars = len(text)

        # Split into sections by headers (markdown-style)
        lines = text.split("\n")
        current_section_title = "Content"
        current_section_content = []

        for line in lines:
            # Check for markdown headers
            if line.startswith("#"):
                if current_section_content:
                    section = self._create_section(
                        title=current_section_title,
                        content="\n".join(current_section_content),
                    )
                    result.sections.append(section)

                current_section_title = line.lstrip("#").strip()
                current_section_content = []
            else:
                current_section_content.append(line)

        # Save last section
        if current_section_content:
            section = self._create_section(
                title=current_section_title,
                content="\n".join(current_section_content),
            )
            result.sections.append(section)

        result.security_excerpts = self._extract_security_excerpts(text)

        if description:
            result.metadata["user_description"] = description

        return result

    def _create_section(
        self,
        title: str,
        content: str,
        page_number: Optional[int] = None,
        section_type: str = "text",
    ) -> DocumentSection:
        """Create a document section with importance scoring."""
        # Calculate importance based on security keywords
        content_lower = content.lower()
        title_lower = title.lower()
        combined = content_lower + " " + title_lower

        keywords_found = []
        importance_score = 0.3  # Base score

        for keyword in SECURITY_KEYWORDS:
            if keyword in combined:
                keywords_found.append(keyword)
                importance_score += 0.05  # Increase for each keyword

        # Cap at 1.0
        importance_score = min(1.0, importance_score)

        return DocumentSection(
            title=title,
            content=content,
            page_number=page_number,
            section_type=section_type,
            importance_score=importance_score,
            keywords_found=keywords_found,
        )

    def _extract_security_excerpts(
        self,
        text: str,
        max_excerpts: int = 20,
        context_chars: int = 500,
    ) -> List[str]:
        """Extract security-relevant excerpts from text."""
        if not text:
            return []

        text_lower = text.lower()
        excerpts = []
        used_positions = set()

        for keyword in SECURITY_KEYWORDS:
            # Find all occurrences
            start = 0
            while True:
                pos = text_lower.find(keyword, start)
                if pos == -1:
                    break

                # Check if we've already used this area
                overlap = False
                for used_pos in used_positions:
                    if abs(pos - used_pos) < context_chars:
                        overlap = True
                        break

                if not overlap:
                    # Extract context around the keyword
                    excerpt_start = max(0, pos - context_chars // 2)
                    excerpt_end = min(len(text), pos + len(keyword) + context_chars // 2)

                    excerpt = text[excerpt_start:excerpt_end].strip()
                    if excerpt_start > 0:
                        excerpt = "..." + excerpt
                    if excerpt_end < len(text):
                        excerpt = excerpt + "..."

                    excerpts.append(excerpt)
                    used_positions.add(pos)

                    if len(excerpts) >= max_excerpts:
                        return excerpts

                start = pos + 1

        return excerpts

    def parse_documents_batch(
        self,
        documents: List[Dict[str, Any]],
    ) -> List[ParsedDocument]:
        """Parse multiple documents."""
        results = []
        for doc in documents:
            parsed = self.parse_document(
                filename=doc.get("filename", "unknown"),
                content_base64=doc.get("content_base64", ""),
                content_type=doc.get("content_type"),
                description=doc.get("description"),
                ocr_languages=doc.get("ocr_languages"),
            )
            results.append(parsed)
        return results

    def correlate_findings_with_documents(
        self,
        findings: List[Dict[str, Any]],
        parsed_documents: List[ParsedDocument],
        max_correlations_per_finding: int = 5,
    ) -> List[Dict[str, Any]]:
        """
        Correlate security findings with relevant document sections.

        For each finding, searches parsed documents for content that mentions:
        - The affected endpoint/path
        - The vulnerability type
        - Parameter names
        - File paths
        - Function names

        Args:
            findings: List of security findings with fields like:
                - endpoint/path/url
                - vulnerability_type
                - parameter
                - file_path
                - description
            parsed_documents: List of ParsedDocument objects
            max_correlations_per_finding: Max document references per finding

        Returns:
            List of findings enriched with 'document_correlations' field
        """
        if not parsed_documents:
            return findings

        enriched_findings = []

        for finding in findings:
            correlations = []
            search_terms = self._extract_search_terms_from_finding(finding)

            if not search_terms:
                enriched_findings.append(finding)
                continue

            # Search each document for relevant content
            for doc in parsed_documents:
                doc_correlations = self._search_document_for_terms(
                    doc, search_terms, max_correlations_per_finding
                )
                correlations.extend(doc_correlations)

            # Sort by relevance score and limit
            correlations.sort(key=lambda x: x.get("relevance_score", 0), reverse=True)
            correlations = correlations[:max_correlations_per_finding]

            # Add correlations to finding
            enriched = dict(finding)
            if correlations:
                enriched["document_correlations"] = correlations
                enriched["has_documentation"] = True
            else:
                enriched["document_correlations"] = []
                enriched["has_documentation"] = False

            enriched_findings.append(enriched)

        return enriched_findings

    def _extract_search_terms_from_finding(
        self,
        finding: Dict[str, Any],
    ) -> List[Tuple[str, float]]:
        """
        Extract searchable terms from a finding with priority weights.

        Returns list of (term, weight) tuples where weight indicates importance.
        """
        terms = []

        # Endpoint/path - highest priority
        endpoint = finding.get("endpoint") or finding.get("path") or finding.get("url")
        if endpoint:
            # Extract path from full URL if needed
            if "://" in endpoint:
                try:
                    from urllib.parse import urlparse
                    parsed = urlparse(endpoint)
                    endpoint = parsed.path
                except Exception:
                    pass

            terms.append((endpoint, 1.0))

            # Also add path segments
            segments = [s for s in endpoint.split("/") if s and len(s) > 2]
            for segment in segments:
                if not segment.startswith("{"):  # Skip path params like {id}
                    terms.append((segment, 0.7))

        # Parameter names - high priority
        param = finding.get("parameter") or finding.get("param_name")
        if param:
            terms.append((param, 0.9))

        # File path
        file_path = finding.get("file_path") or finding.get("file")
        if file_path:
            terms.append((file_path, 0.9))
            # Also add filename
            filename = os.path.basename(file_path)
            if filename:
                terms.append((filename, 0.7))

        # Function/method name
        function = finding.get("function") or finding.get("method_name")
        if function:
            terms.append((function, 0.8))

        # Vulnerability type keywords
        vuln_type = finding.get("vulnerability_type") or finding.get("type") or finding.get("category")
        if vuln_type:
            # Map common vuln types to search terms
            vuln_terms = self._get_vuln_type_search_terms(vuln_type)
            for vt in vuln_terms:
                terms.append((vt, 0.5))

        # Line number context (for code findings)
        line_num = finding.get("line_number") or finding.get("line")
        if line_num and file_path:
            # Look for code snippets mentioning this location
            terms.append((f"line {line_num}", 0.4))

        # CWE/CVE references
        cwe = finding.get("cwe") or finding.get("cwe_id")
        if cwe:
            terms.append((str(cwe), 0.6))

        cve = finding.get("cve") or finding.get("cve_id")
        if cve:
            terms.append((str(cve), 0.8))

        return terms

    def _get_vuln_type_search_terms(self, vuln_type: str) -> List[str]:
        """Map vulnerability types to relevant search terms."""
        vuln_type_lower = vuln_type.lower()

        mappings = {
            "sql_injection": ["sql", "query", "database", "injection", "prepared statement"],
            "sqli": ["sql", "query", "database", "injection"],
            "xss": ["xss", "cross-site", "script", "sanitize", "escape", "html"],
            "cross_site_scripting": ["xss", "cross-site", "script", "sanitize"],
            "csrf": ["csrf", "cross-site request", "token", "origin"],
            "ssrf": ["ssrf", "server-side request", "url", "fetch", "request"],
            "idor": ["idor", "authorization", "access control", "object reference"],
            "path_traversal": ["path", "traversal", "directory", "file", "../"],
            "rce": ["rce", "command", "execution", "shell", "system"],
            "command_injection": ["command", "injection", "shell", "exec", "system"],
            "xxe": ["xxe", "xml", "entity", "parser"],
            "deserialization": ["deserialize", "serialize", "pickle", "unmarshall"],
            "auth_bypass": ["authentication", "bypass", "auth", "login"],
            "sensitive_data": ["sensitive", "exposure", "data", "leak", "pii"],
            "file_upload": ["upload", "file", "extension", "mime"],
        }

        for key, terms in mappings.items():
            if key in vuln_type_lower:
                return terms

        # Default: return the vuln type itself
        return [vuln_type_lower]

    def _search_document_for_terms(
        self,
        doc: ParsedDocument,
        search_terms: List[Tuple[str, float]],
        max_results: int = 5,
    ) -> List[Dict[str, Any]]:
        """
        Search a parsed document for terms and return relevant excerpts.
        """
        results = []
        full_text_lower = doc.full_text.lower() if doc.full_text else ""

        # Track which positions we've already extracted to avoid duplicates
        used_positions = set()

        for term, weight in search_terms:
            term_lower = term.lower()

            # Skip very short terms
            if len(term_lower) < 3:
                continue

            # Search in full text
            pos = 0
            while pos < len(full_text_lower):
                found_pos = full_text_lower.find(term_lower, pos)
                if found_pos == -1:
                    break

                # Check if we've already used this area
                is_duplicate = any(abs(found_pos - used) < 300 for used in used_positions)
                if not is_duplicate:
                    # Extract context around the match
                    context_start = max(0, found_pos - 200)
                    context_end = min(len(doc.full_text), found_pos + len(term) + 200)

                    excerpt = doc.full_text[context_start:context_end].strip()
                    if context_start > 0:
                        excerpt = "..." + excerpt
                    if context_end < len(doc.full_text):
                        excerpt = excerpt + "..."

                    # Calculate relevance score
                    relevance_score = weight

                    # Boost if term appears in high-importance section
                    for section in doc.sections:
                        if term_lower in section.content.lower():
                            relevance_score += section.importance_score * 0.3
                            break

                    results.append({
                        "document": doc.filename,
                        "document_type": doc.document_type.value,
                        "matched_term": term,
                        "excerpt": excerpt,
                        "relevance_score": min(1.0, relevance_score),
                        "position": found_pos,
                    })

                    used_positions.add(found_pos)

                pos = found_pos + 1

                if len(results) >= max_results * 2:  # Get extras to sort later
                    break

        # Also check API endpoints for path matches
        for endpoint in doc.api_endpoints:
            endpoint_path = endpoint.get("path", "")
            for term, weight in search_terms:
                if term.lower() in endpoint_path.lower():
                    results.append({
                        "document": doc.filename,
                        "document_type": "api_endpoint",
                        "matched_term": term,
                        "excerpt": f"API: {endpoint.get('method', 'GET')} {endpoint_path} - {endpoint.get('summary', '')} {endpoint.get('description', '')}",
                        "relevance_score": min(1.0, weight + 0.2),
                        "endpoint_details": endpoint,
                    })

        # Sort by relevance and return top results
        results.sort(key=lambda x: x.get("relevance_score", 0), reverse=True)
        return results[:max_results]

    def build_finding_context_from_documents(
        self,
        finding: Dict[str, Any],
        parsed_documents: List[ParsedDocument],
        max_context_chars: int = 2000,
    ) -> str:
        """
        Build a context string for a specific finding from related documents.

        This is useful for providing AI agents with relevant documentation
        context when analyzing a specific finding.
        """
        # Get correlations
        correlations = self.correlate_findings_with_documents(
            [finding], parsed_documents, max_correlations_per_finding=10
        )[0].get("document_correlations", [])

        if not correlations:
            return ""

        context_parts = ["## Related Documentation\n"]
        current_chars = 0

        for corr in correlations:
            entry = f"\n**From {corr['document']} (matched: {corr['matched_term']}):**\n{corr['excerpt']}\n"

            if current_chars + len(entry) > max_context_chars:
                break

            context_parts.append(entry)
            current_chars += len(entry)

        return "".join(context_parts)

    def get_document_summary_for_findings(
        self,
        findings: List[Dict[str, Any]],
        parsed_documents: List[ParsedDocument],
    ) -> Dict[str, Any]:
        """
        Generate a summary of how documents relate to findings.

        Returns statistics and high-level correlations useful for
        understanding documentation coverage.
        """
        # Correlate all findings
        enriched = self.correlate_findings_with_documents(findings, parsed_documents)

        # Calculate statistics
        findings_with_docs = sum(1 for f in enriched if f.get("has_documentation"))
        findings_without_docs = len(enriched) - findings_with_docs

        # Group by document
        docs_referenced = {}
        for finding in enriched:
            for corr in finding.get("document_correlations", []):
                doc_name = corr["document"]
                if doc_name not in docs_referenced:
                    docs_referenced[doc_name] = {
                        "finding_count": 0,
                        "matched_terms": set(),
                    }
                docs_referenced[doc_name]["finding_count"] += 1
                docs_referenced[doc_name]["matched_terms"].add(corr["matched_term"])

        # Convert sets to lists for JSON serialization
        for doc_info in docs_referenced.values():
            doc_info["matched_terms"] = list(doc_info["matched_terms"])

        # Find undocumented endpoints
        undocumented_endpoints = []
        for finding in enriched:
            if not finding.get("has_documentation"):
                endpoint = finding.get("endpoint") or finding.get("path")
                if endpoint:
                    undocumented_endpoints.append(endpoint)

        return {
            "total_findings": len(findings),
            "findings_with_documentation": findings_with_docs,
            "findings_without_documentation": findings_without_docs,
            "documentation_coverage_percent": round(
                (findings_with_docs / len(findings) * 100) if findings else 0, 1
            ),
            "documents_referenced": docs_referenced,
            "undocumented_endpoints": undocumented_endpoints[:20],  # Top 20
            "total_documents": len(parsed_documents),
        }


# Singleton instance
document_parser_service = DocumentParserService()
