"""
Document AI Service - Extract text from documents and analyze with Gemini AI.

Supports: PDF, Word (docx), PowerPoint (pptx), Text files, and more.
"""

import os
import io
import json
import logging
import shutil
from pathlib import Path
from typing import Optional, List, Dict, Any
import asyncio

# Document extraction libraries
try:
    import PyPDF2
    HAS_PYPDF = True
except ImportError:
    HAS_PYPDF = False

try:
    import fitz  # PyMuPDF
    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False

try:
    import pytesseract
    from PIL import Image
    HAS_PYTESSERACT = True
except ImportError:
    HAS_PYTESSERACT = False
    pytesseract = None
    Image = None

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# AI libraries
try:
    from google import genai
    HAS_GENAI = True
except ImportError:
    HAS_GENAI = False
    genai = None
from backend.core.config import settings

logger = logging.getLogger(__name__)


class DocumentAIService:
    """Service for document text extraction and AI analysis."""
    
    SUPPORTED_MIME_TYPES = {
        'application/pdf': 'pdf',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
        'application/msword': 'doc',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
        'application/vnd.ms-powerpoint': 'ppt',
        'text/plain': 'txt',
        'text/markdown': 'md',
        'text/csv': 'csv',
        'application/json': 'json',
        'application/xml': 'xml',
        'text/xml': 'xml',
    }
    
    MAX_TEXT_LENGTH = 500000  # Max characters for standard AI analysis
    MAX_DEEP_TEXT_LENGTH = 2000000  # Max characters for deep analysis
    CHUNK_SIZE = 4000  # Characters per chunk for standard context
    DEEP_CHUNK_SIZE = 8000  # Characters per chunk for deep analysis
    MAX_DEEP_CHUNKS = 60  # Max chunks to process in deep mode
    
    def __init__(self):
        """Initialize the document AI service."""
        self.gemini_api_key = settings.gemini_api_key
        if self.gemini_api_key and HAS_GENAI:
            self.client = genai.Client(api_key=self.gemini_api_key)
            self.model_name = settings.gemini_model_id
        else:
            self.client = None
            self.model_name = None
            if not HAS_GENAI:
                logger.warning("google-genai not installed - document AI features disabled")
            else:
                logger.warning("GEMINI_API_KEY not set - document AI features disabled")
        self._has_pymupdf = HAS_PYMUPDF
        self._has_pytesseract = HAS_PYTESSERACT
        self._has_tesseract = self._has_pytesseract and shutil.which("tesseract") is not None
        if not self._has_tesseract and self._has_pytesseract:
            logger.info("pytesseract installed but tesseract binary not found in PATH")
    
    def is_supported(self, mime_type: str) -> bool:
        """Check if a file type is supported for text extraction."""
        return mime_type in self.SUPPORTED_MIME_TYPES
    
    def extract_text(self, file_path: str, mime_type: str, ocr_languages: Optional[str] = None) -> str:
        """
        Extract text content from a document file.
        
        Args:
            file_path: Path to the file
            mime_type: MIME type of the file
            ocr_languages: Optional tesseract language codes (e.g., "eng", "rus+eng")
            
        Returns:
            Extracted text content
        """
        file_type = self.SUPPORTED_MIME_TYPES.get(mime_type, '')
        
        try:
            if file_type == 'pdf':
                return self._extract_pdf(file_path, ocr_languages)
            elif file_type in ('docx', 'doc'):
                return self._extract_docx(file_path)
            elif file_type in ('pptx', 'ppt'):
                return self._extract_pptx(file_path)
            elif file_type in ('txt', 'md', 'csv', 'json', 'xml'):
                return self._extract_text_file(file_path)
            else:
                raise ValueError(f"Unsupported file type: {mime_type}")
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {e}")
            raise
    
    def _extract_pdf(self, file_path: str, ocr_languages: Optional[str] = None) -> str:
        """Extract text from PDF file."""
        text_parts = []
        ocr_langs = ocr_languages or "eng"

        if self._has_pymupdf:
            try:
                pdf = fitz.open(file_path)
                for page_num, page in enumerate(pdf):
                    page_text = page.get_text() or ""
                    if not page_text.strip() and self._has_tesseract:
                        page_text = self._ocr_pdf_page(page, ocr_langs)
                    if page_text:
                        text_parts.append(f"--- Page {page_num + 1} ---\n{page_text}")
                pdf.close()
                if text_parts:
                    return "\n\n".join(text_parts)
            except Exception as e:
                logger.warning(f"PyMuPDF extraction failed, falling back to PyPDF2: {e}")

        if not HAS_PYPDF:
            raise ImportError("PyPDF2 not installed. Install with: pip install PyPDF2")
        
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page_num, page in enumerate(reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(f"--- Page {page_num + 1} ---\n{page_text}")
                except Exception as e:
                    logger.warning(f"Error extracting page {page_num + 1}: {e}")
        
        return "\n\n".join(text_parts)

    def _ocr_pdf_page(self, page: "fitz.Page", ocr_languages: str) -> str:
        """Run OCR on a PyMuPDF page if tesseract is available."""
        if not self._has_tesseract or not self._has_pytesseract:
            return ""
        try:
            pix = page.get_pixmap(dpi=300)
            mode = "RGB" if pix.alpha == 0 else "RGBA"
            image = Image.frombytes(mode, [pix.width, pix.height], pix.samples)
            if mode == "RGBA":
                image = image.convert("RGB")
            config = "--oem 3 --psm 6"
            return pytesseract.image_to_string(image, lang=ocr_languages, config=config) or ""
        except Exception as e:
            logger.warning(f"OCR failed for PDF page: {e}")
            return ""
    
    def _extract_docx(self, file_path: str) -> str:
        """Extract text from Word document."""
        if not HAS_DOCX:
            raise ImportError("python-docx not installed. Install with: pip install python-docx")
        
        doc = DocxDocument(file_path)
        text_parts = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)
        
        return "\n\n".join(text_parts)
    
    def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint presentation."""
        if not HAS_PPTX:
            raise ImportError("python-pptx not installed. Install with: pip install python-pptx")
        
        prs = Presentation(file_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text = [f"--- Slide {slide_num + 1} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if len(slide_text) > 1:
                text_parts.append("\n".join(slide_text))
        
        return "\n\n".join(text_parts)
    
    def _extract_text_file(self, file_path: str) -> str:
        """Extract text from plain text files."""
        encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    return f.read()
            except (UnicodeDecodeError, UnicodeError):
                continue
        
        # Fallback: read as binary and decode with errors='replace'
        with open(file_path, 'rb') as f:
            return f.read().decode('utf-8', errors='replace')
    
    def chunk_text(self, text: str, chunk_size: int = None) -> List[str]:
        """
        Split text into chunks for processing.
        
        Args:
            text: Full text content
            chunk_size: Max characters per chunk
            
        Returns:
            List of text chunks
        """
        chunk_size = chunk_size or self.CHUNK_SIZE
        chunks = []
        
        # Split by paragraphs first
        paragraphs = text.split('\n\n')
        current_chunk = []
        current_length = 0
        
        for para in paragraphs:
            para_length = len(para)
            
            if current_length + para_length > chunk_size and current_chunk:
                chunks.append('\n\n'.join(current_chunk))
                current_chunk = []
                current_length = 0
            
            # If a single paragraph is too long, split it
            if para_length > chunk_size:
                if current_chunk:
                    chunks.append('\n\n'.join(current_chunk))
                    current_chunk = []
                    current_length = 0
                
                # Split long paragraph by sentences
                sentences = para.replace('. ', '.\n').split('\n')
                for sentence in sentences:
                    if len(sentence) > chunk_size:
                        # Split by words as last resort
                        words = sentence.split()
                        temp_chunk = []
                        temp_length = 0
                        for word in words:
                            if temp_length + len(word) + 1 > chunk_size:
                                if temp_chunk:
                                    chunks.append(' '.join(temp_chunk))
                                temp_chunk = [word]
                                temp_length = len(word)
                            else:
                                temp_chunk.append(word)
                                temp_length += len(word) + 1
                        if temp_chunk:
                            current_chunk = [' '.join(temp_chunk)]
                            current_length = temp_length
                    else:
                        current_chunk.append(sentence)
                        current_length += len(sentence) + 2
            else:
                current_chunk.append(para)
                current_length += para_length + 2
        
        if current_chunk:
            chunks.append('\n\n'.join(current_chunk))
        
        return chunks
    
    async def generate_summary(
        self,
        text: str,
        filename: str,
        analysis_depth: str = "standard",
    ) -> Dict[str, Any]:
        """
        Generate AI summary of document content.
        
        Args:
            text: Extracted text content
            filename: Original filename for context
            analysis_depth: "standard" or "deep"
            
        Returns:
            Dict with summary and key_points
        """
        if not self.client:
            raise ValueError("Gemini API key not configured or google-genai not installed")

        if analysis_depth.lower() == "deep":
            return await self._generate_deep_summary(text, filename)
        
        # Truncate if too long
        if len(text) > self.MAX_TEXT_LENGTH:
            text = text[:self.MAX_TEXT_LENGTH] + "\n\n[Content truncated due to length...]"
        
        prompt = f"""Analyze the following document and provide a comprehensive analysis.

**IMPORTANT INSTRUCTIONS:**
- If the document is in a language other than English, TRANSLATE all content to English in your response
- ONLY include information that is explicitly stated in the document - do NOT invent, assume, or hallucinate any details
- If something is unclear or missing from the document, state that explicitly rather than guessing
- Use rich formatting in your response for readability:
  • Use **bold** for key terms, names, and important concepts
  • Use *italics* for emphasis or technical terms
  • Use bullet points (•) for lists
  • Use numbered lists where sequence matters
  • Break long text into clear paragraphs with headers if needed

Document name: {filename}

---
{text}
---

Provide your analysis in the following JSON format:
{{
    "summary": "Your comprehensive summary here (2-4 paragraphs, well-formatted with **bold** and *italics* for readability, use • for bullet points within text if helpful)...",
    "key_points": [
        "• **Key Point 1** - Brief explanation with relevant details",
        "• **Key Point 2** - Brief explanation with relevant details",
        "• Continue with 5-10 important takeaways..."
    ],
    "language_detected": "Original language of the document (e.g., English, Spanish, Chinese, etc.)",
    "was_translated": true/false
}}

Remember: Accuracy is paramount. Only report what the document actually says.
"""
        
        try:
            from google.genai import types
            response = await asyncio.to_thread(
                self.client.models.generate_content,
                model=self.model_name,
                contents=prompt,
                config=types.GenerateContentConfig(
                    thinking_config=types.ThinkingConfig(thinking_level="medium"),
                    max_output_tokens=2000,
                )
            )
            
            # Parse JSON response
            response_text = response.text.strip()
            
            # Try to extract JSON from response
            if response_text.startswith('```json'):
                response_text = response_text[7:]
            if response_text.startswith('```'):
                response_text = response_text[3:]
            if response_text.endswith('```'):
                response_text = response_text[:-3]
            
            result = json.loads(response_text.strip())
            return {
                "summary": result.get("summary", ""),
                "key_points": result.get("key_points", []),
                "language_detected": result.get("language_detected", "Unknown"),
                "was_translated": result.get("was_translated", False)
            }
            
        except json.JSONDecodeError:
            # If JSON parsing fails, return the raw text as summary
            return {
                "summary": response.text if response else "Failed to generate summary",
                "key_points": [],
                "language_detected": "Unknown",
                "was_translated": False
            }
        except Exception as e:
            logger.error(f"Error generating summary: {e}")
            raise

    async def _generate_deep_summary(self, text: str, filename: str) -> Dict[str, Any]:
        """Generate a detailed summary using chunked analysis."""
        if len(text) > self.MAX_DEEP_TEXT_LENGTH:
            text = text[:self.MAX_DEEP_TEXT_LENGTH] + "\n\n[Content truncated due to length...]"
            truncated = True
        else:
            truncated = False

        chunks = self.chunk_text(text, chunk_size=self.DEEP_CHUNK_SIZE)
        if len(chunks) > self.MAX_DEEP_CHUNKS:
            chunks = chunks[:self.MAX_DEEP_CHUNKS]
            truncated = True

        chunk_summaries = []
        all_key_points: List[str] = []
        detected_language = "Unknown"
        was_translated = False

        for idx, chunk in enumerate(chunks, 1):
            result = await self._summarize_chunk(chunk, filename, idx, len(chunks))
            chunk_summary = result.get("chunk_summary", "").strip()
            if chunk_summary:
                chunk_summaries.append(f"### Section {idx}\n{chunk_summary}")
            key_points = result.get("key_points", [])
            if isinstance(key_points, list):
                all_key_points.extend(key_points)
            if detected_language == "Unknown" and result.get("language_detected"):
                detected_language = result.get("language_detected") or detected_language
            if result.get("was_translated"):
                was_translated = True

        exec_summary = await self._synthesize_chunk_summaries(chunk_summaries, filename)
        detailed_summary = "\n\n".join(chunk_summaries)

        summary_parts = []
        if exec_summary:
            summary_parts.append("## Executive Summary\n" + exec_summary)
        summary_parts.append("## Detailed Analysis\n" + detailed_summary)
        if truncated:
            summary_parts.append("\n\n[Content truncated due to length...]")

        summary_text = "\n\n".join(summary_parts).strip()
        key_points = self._dedupe_key_points(all_key_points, limit=40)

        return {
            "summary": summary_text,
            "key_points": key_points,
            "language_detected": detected_language,
            "was_translated": was_translated,
        }

    async def _summarize_chunk(
        self,
        chunk: str,
        filename: str,
        chunk_index: int,
        total_chunks: int,
    ) -> Dict[str, Any]:
        """Summarize a document chunk with strict, extractive guidance."""
        prompt = f"""Analyze this chunk from a document and provide an extractive summary.

**CRITICAL INSTRUCTIONS:**
- ONLY include information explicitly stated in the chunk
- Do NOT invent, assume, or hallucinate
- If the chunk is not in English, translate its content to English
- Keep wording close to the source where possible (extractive summary)

Document: {filename}
Chunk: {chunk_index} of {total_chunks}

---
{chunk}
---

Provide JSON only:
{{
  "chunk_summary": "1-2 paragraphs summarizing only this chunk",
  "key_points": [
    "• Key point 1",
    "• Key point 2",
    "• 5-10 points total"
  ],
  "language_detected": "Language name",
  "was_translated": true/false
}}
"""
        try:
            from google.genai import types
            response = await asyncio.to_thread(
                self.client.models.generate_content,
                model=self.model_name,
                contents=prompt,
                config=types.GenerateContentConfig(
                    thinking_config=types.ThinkingConfig(thinking_level="low"),
                    max_output_tokens=1500,
                )
            )
            return self._parse_json_response(response.text)
        except Exception as e:
            logger.error(f"Chunk summary failed: {e}")
            return {"chunk_summary": "", "key_points": []}

    async def _synthesize_chunk_summaries(self, chunk_summaries: List[str], filename: str) -> str:
        """Generate an executive summary from chunk summaries."""
        if not chunk_summaries:
            return ""
        summary_text = "\n\n".join(chunk_summaries)
        if len(summary_text) > self.MAX_TEXT_LENGTH:
            summary_text = summary_text[:self.MAX_TEXT_LENGTH] + "\n\n[Content truncated due to length...]"

        prompt = f"""Create an executive summary using ONLY the provided chunk summaries.

**CRITICAL INSTRUCTIONS:**
- Do NOT invent or add details not present
- Summarize in 2-3 paragraphs
- Keep the summary faithful to the source

Document: {filename}

---
{summary_text}
---

Return JSON only:
{{
  "summary": "Executive summary here",
  "key_points": [
    "• Key point 1",
    "• Key point 2",
    "• 8-12 points total"
  ]
}}
"""
        try:
            from google.genai import types
            response = await asyncio.to_thread(
                self.client.models.generate_content,
                model=self.model_name,
                contents=prompt,
                config=types.GenerateContentConfig(
                    thinking_config=types.ThinkingConfig(thinking_level="low"),
                    max_output_tokens=1200,
                )
            )
            result = self._parse_json_response(response.text)
            return result.get("summary", "").strip()
        except Exception as e:
            logger.error(f"Executive summary synthesis failed: {e}")
            return ""

    def _parse_json_response(self, response_text: str) -> Dict[str, Any]:
        """Parse JSON from a model response, handling code fences."""
        if not response_text:
            return {}
        text = response_text.strip()
        if text.startswith('```json'):
            text = text[7:]
        if text.startswith('```'):
            text = text[3:]
        if text.endswith('```'):
            text = text[:-3]
        try:
            return json.loads(text.strip())
        except json.JSONDecodeError:
            return {}

    def _dedupe_key_points(self, points: List[str], limit: int = 40) -> List[str]:
        """Deduplicate key points while preserving order."""
        seen = set()
        deduped = []
        for point in points:
            clean = point.strip()
            if not clean:
                continue
            key = clean.lower()
            if key in seen:
                continue
            seen.add(key)
            deduped.append(clean)
            if len(deduped) >= limit:
                break
        return deduped
    
    async def generate_multi_document_summary(
        self,
        documents: List[Dict[str, str]],
        custom_prompt: str = "",
        analysis_depth: str = "standard",
    ) -> Dict[str, Any]:
        """
        Generate AI summary across multiple documents.
        
        Args:
            documents: List of dicts with 'filename' and 'text' keys
            custom_prompt: User-provided additional instructions
            analysis_depth: "standard" or "deep"
            
        Returns:
            Dict with combined_summary and combined_key_points
        """
        if not self.client:
            raise ValueError("Gemini API key not configured or google-genai not installed")

        if analysis_depth.lower() == "deep":
            return await self._generate_deep_multi_document_summary(documents, custom_prompt)
        
        # Build combined document text
        combined_parts = []
        for i, doc in enumerate(documents, 1):
            text = doc.get("text", "")
            # Truncate individual docs if needed
            if len(text) > 100000:
                text = text[:100000] + "\n[Content truncated due to length...]"
            combined_parts.append(f"=== DOCUMENT {i}: {doc.get('filename', 'Unknown')} ===\n{text}")
        
        combined_text = "\n\n".join(combined_parts)
        
        # Further truncate if total is too long
        if len(combined_text) > self.MAX_TEXT_LENGTH:
            combined_text = combined_text[:self.MAX_TEXT_LENGTH] + "\n\n[Content truncated due to length...]"
        
        # Build the prompt
        custom_instructions = ""
        if custom_prompt and custom_prompt.strip():
            custom_instructions = f"""
**USER-PROVIDED INSTRUCTIONS (PRIORITY):**
{custom_prompt.strip()}

"""
        
        prompt = f"""Analyze the following {len(documents)} document(s) and provide a comprehensive combined analysis.

{custom_instructions}**IMPORTANT INSTRUCTIONS:**
- If any document is in a language other than English, TRANSLATE all content to English in your response
- ONLY include information that is explicitly stated in the documents - do NOT invent, assume, or hallucinate any details
- If something is unclear or missing from the documents, state that explicitly rather than guessing
- When analyzing multiple documents, identify connections, patterns, and relationships between them
- Use rich formatting in your response for readability:
  • Use **bold** for key terms, names, and important concepts
  • Use *italics* for emphasis or technical terms
  • Use bullet points (•) for lists
  • Use numbered lists where sequence matters
  • Reference specific documents by name when relevant

---
{combined_text}
---

Provide your analysis in the following JSON format:
{{
    "combined_summary": "Your comprehensive combined summary here (3-5 paragraphs, well-formatted with **bold** and *italics*, reference documents by name, identify cross-document patterns)...",
    "combined_key_points": [
        "• **Key Finding 1** - Brief explanation with relevant details (cite which document(s))",
        "• **Key Finding 2** - Brief explanation with relevant details",
        "• Continue with 8-15 important takeaways across all documents..."
    ],
    "documents_analyzed": {len(documents)},
    "languages_detected": ["List of languages found in the documents"]
}}

Remember: Accuracy is paramount. Only report what the documents actually say.
"""
        
        try:
            from google.genai import types
            response = await asyncio.to_thread(
                self.client.models.generate_content,
                model=self.model_name,
                contents=prompt,
                config=types.GenerateContentConfig(
                    thinking_config=types.ThinkingConfig(thinking_level="medium"),
                    max_output_tokens=4000,
                )
            )
            
            # Parse JSON response
            response_text = response.text.strip()
            
            # Try to extract JSON from response
            if response_text.startswith('```json'):
                response_text = response_text[7:]
            if response_text.startswith('```'):
                response_text = response_text[3:]
            if response_text.endswith('```'):
                response_text = response_text[:-3]
            
            result = json.loads(response_text.strip())
            return {
                "combined_summary": result.get("combined_summary", ""),
                "combined_key_points": result.get("combined_key_points", []),
                "documents_analyzed": result.get("documents_analyzed", len(documents)),
                "languages_detected": result.get("languages_detected", [])
            }
            
        except json.JSONDecodeError:
            return {
                "combined_summary": response.text if response else "Failed to generate summary",
                "combined_key_points": [],
                "documents_analyzed": len(documents),
                "languages_detected": []
            }
        except Exception as e:
            logger.error(f"Error generating multi-document summary: {e}")
            raise

    async def _generate_deep_multi_document_summary(
        self,
        documents: List[Dict[str, str]],
        custom_prompt: str = "",
    ) -> Dict[str, Any]:
        """Generate a detailed combined summary using deep per-document analysis."""
        doc_summaries = []
        all_key_points: List[str] = []
        languages_detected = set()

        for doc in documents:
            filename = doc.get("filename", "Unknown")
            text = doc.get("text", "")
            result = await self._generate_deep_summary(text, filename)
            summary = result.get("summary", "").strip()
            if summary:
                doc_summaries.append(f"### {filename}\n{summary}")
            key_points = result.get("key_points", [])
            if isinstance(key_points, list):
                all_key_points.extend(key_points)
            language = result.get("language_detected")
            if language:
                languages_detected.add(language)

        combined_exec_summary = ""
        if doc_summaries:
            combined_exec_summary = await self._synthesize_chunk_summaries(doc_summaries, "Combined Documents")

        combined_summary_parts = []
        if combined_exec_summary:
            combined_summary_parts.append("## Executive Summary\n" + combined_exec_summary)

        if custom_prompt and custom_prompt.strip():
            combined_summary_parts.append(
                f"## User Requirements\n{custom_prompt.strip()}"
            )

        combined_summary_parts.append("## Detailed Document Analysis\n" + "\n\n".join(doc_summaries))
        combined_summary = "\n\n".join(combined_summary_parts).strip()

        return {
            "combined_summary": combined_summary,
            "combined_key_points": self._dedupe_key_points(all_key_points, limit=60),
            "documents_analyzed": len(documents),
            "languages_detected": sorted(languages_detected),
        }
    
    async def answer_question(
        self, 
        question: str, 
        document_text: str, 
        document_summary: str,
        chat_history: List[Dict[str, str]] = None,
        filename: str = ""
    ) -> str:
        """
        Answer a question about the document using AI.
        
        Args:
            question: User's question
            document_text: Full document text (truncated if needed)
            document_summary: Previously generated summary
            chat_history: Previous Q&A in this conversation
            filename: Document filename for context
            
        Returns:
            AI-generated answer
        """
        if not self.client:
            raise ValueError("Gemini API key not configured or google-genai not installed")
        
        # Build context with summary and relevant text chunks
        context_text = document_text
        if len(context_text) > 50000:
            context_text = context_text[:50000] + "\n\n[Content truncated...]"
        
        # Build chat history string
        history_str = ""
        if chat_history:
            for msg in chat_history[-10:]:  # Last 10 messages for context
                role = "User" if msg["role"] == "user" else "Assistant"
                history_str += f"\n{role}: {msg['content']}\n"
        
        prompt = f"""You are a helpful assistant answering questions about a document.

**CRITICAL INSTRUCTIONS:**
- If the document is in a foreign language, provide your answer in ENGLISH (translate relevant content)
- ONLY answer based on what is explicitly stated in the document
- Do NOT invent, assume, or hallucinate any information
- If the answer is not in the document, clearly state: "This information is not found in the document."
- If you're uncertain, express that uncertainty rather than guessing

**FORMATTING:**
- Use **bold** for key terms, names, and important concepts
- Use *italics* for emphasis or technical terms
- Use bullet points (•) for lists when appropriate
- Keep your answer clear and well-structured

Document: {filename}

Document Summary:
{document_summary}

Full Document Content:
---
{context_text}
---

{f"Previous conversation:{history_str}" if history_str else ""}

User Question: {question}

Provide a helpful, accurate, well-formatted answer based ONLY on the document content."""

        try:
            from google.genai import types
            response = await asyncio.to_thread(
                self.client.models.generate_content,
                model=self.model_name,
                contents=prompt,
                config=types.GenerateContentConfig(
                    thinking_config=types.ThinkingConfig(thinking_level="high"),
                    max_output_tokens=1500,
                )
            )
            return response.text.strip()
        except Exception as e:
            logger.error(f"Error answering question: {e}")
            raise
    
    async def answer_report_question(
        self, 
        question: str, 
        documents: List[Dict[str, str]],
        report_summary: str,
        chat_history: List[Dict[str, str]] = None,
        custom_prompt: str = ""
    ) -> str:
        """
        Answer a question about multiple documents in a report.
        
        Args:
            question: User's question
            documents: List of dicts with 'filename' and 'text' keys
            report_summary: Previously generated combined summary
            chat_history: Previous Q&A in this conversation
            custom_prompt: Original custom prompt used for analysis
            
        Returns:
            AI-generated answer
        """
        if not self.client:
            raise ValueError("Gemini API key not configured or google-genai not installed")
        
        # Build combined document context
        combined_parts = []
        for i, doc in enumerate(documents, 1):
            text = doc.get("text", "")[:30000]  # Truncate each doc
            combined_parts.append(f"=== DOCUMENT {i}: {doc.get('filename', 'Unknown')} ===\n{text}")
        
        context_text = "\n\n".join(combined_parts)
        if len(context_text) > 80000:
            context_text = context_text[:80000] + "\n\n[Content truncated...]"
        
        # Build chat history string
        history_str = ""
        if chat_history:
            for msg in chat_history[-10:]:
                role = "User" if msg["role"] == "user" else "Assistant"
                history_str += f"\n{role}: {msg['content']}\n"
        
        custom_instructions = ""
        if custom_prompt and custom_prompt.strip():
            custom_instructions = f"""
**Original Analysis Instructions:**
{custom_prompt.strip()}

"""
        
        prompt = f"""You are a helpful assistant answering questions about a set of {len(documents)} document(s).

{custom_instructions}**CRITICAL INSTRUCTIONS:**
- If documents are in foreign languages, provide your answer in ENGLISH
- ONLY answer based on what is explicitly stated in the documents
- Do NOT invent, assume, or hallucinate any information
- If the answer is not in the documents, clearly state: "This information is not found in the documents."
- Reference specific documents by name when relevant

**FORMATTING:**
- Use **bold** for key terms, names, and important concepts
- Use *italics* for emphasis or technical terms
- Use bullet points (•) for lists when appropriate

Combined Analysis Summary:
{report_summary}

Document Contents:
---
{context_text}
---

{f"Previous conversation:{history_str}" if history_str else ""}

User Question: {question}

Provide a helpful, accurate answer based ONLY on the document content."""

        try:
            from google.genai import types
            response = await asyncio.to_thread(
                self.client.models.generate_content,
                model=self.model_name,
                contents=prompt,
                config=types.GenerateContentConfig(
                    thinking_config=types.ThinkingConfig(thinking_level="high"),
                    max_output_tokens=2000,
                )
            )
            return response.text.strip()
        except Exception as e:
            logger.error(f"Error answering report question: {e}")
            raise


# Singleton instance
document_ai_service = DocumentAIService()
